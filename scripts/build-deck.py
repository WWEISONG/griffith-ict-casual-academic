"""
Build the Head of School briefing deck.

Follows the layout and type conventions of the author's existing Griffith decks:
Arial throughout, a 32pt title at the top-left, a hairline rule beneath it, 17pt
body with square bullets, 10pt sources at the foot, and stat cards carrying a
narrow accent bar down their left edge.

    python3 scripts/build-deck.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# --- palette, taken from the author's existing decks -------------------------
INK    = RGBColor(0x11, 0x11, 0x11)
GREY   = RGBColor(0x59, 0x59, 0x59)
LIGHT  = RGBColor(0xA6, 0xA6, 0xA6)
RED    = RGBColor(0xC0, 0x00, 0x00)
BLUE   = RGBColor(0x08, 0x78, 0xB8)
GREEN  = RGBColor(0x2E, 0x7D, 0x32)
ORANGE = RGBColor(0xC8, 0x74, 0x13)
RULE   = RGBColor(0xD9, 0xD9, 0xD9)
WASH   = RGBColor(0xF5, 0xF5, 0xF5)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

FONT   = "Arial"
T_TITLE = 32
T_LEAD  = 17
T_CARD  = 16
T_BODY  = 12.5
T_NOTE  = 10

# --- the grid ----------------------------------------------------------------
W, H     = Inches(13.333), Inches(7.5)
L        = Inches(0.38)          # left margin
CONTENT  = Inches(12.60)         # content width
RULE_Y   = Inches(1.04)          # hairline under the title
BODY_TOP = Inches(1.92)          # where content begins
FOOT_Y   = Inches(6.89)          # source line
RIGHT_X  = Inches(7.85)          # right-hand column
RIGHT_W  = Inches(5.10)

SHOT  = "docs/tutor-program/screenshots/"
LOGO  = "docs/tutor-program/assets/cap-logo.png"

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def box(s, l, t, w, h):
    tf = s.shapes.add_textbox(l, t, w, h).text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def para(tf, text, size, *, bold=False, color=INK, after=6, before=0,
         align=PP_ALIGN.LEFT, first=False, line=1.15):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.text = text
    p.alignment = align
    p.space_after = Pt(after)
    p.space_before = Pt(before)
    p.line_spacing = line
    f = p.runs[0].font if p.runs else p.font
    f.name, f.size, f.bold, f.color.rgb = FONT, Pt(size), bold, color
    return p


def rect(s, l, t, w, h, fill=None, line=None):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.75)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def head(s, title, lead=None, lead_color=GREY):
    """Title, hairline rule, optional lead line. Returns y for content."""
    tf = box(s, L, Inches(0.08), Inches(9.9), Inches(0.95))
    para(tf, title, T_TITLE, bold=True, first=True, after=0, line=1.08)
    rect(s, 0, RULE_Y, W, Pt(1.1), RULE)
    s.shapes.add_picture(LOGO, Inches(12.55), Inches(0.26), height=Inches(0.62))
    if lead:
        tf = box(s, L, Inches(1.30), CONTENT, Inches(0.5))
        para(tf, lead, T_LEAD, color=lead_color, first=True, after=0, line=1.25)
        return BODY_TOP
    return Inches(1.45)


def bullets(tf, items, size=T_LEAD, color=INK, after=9):
    for i, t in enumerate(items):
        para(tf, "▪  " + t, size, color=color, first=(i == 0), after=after, line=1.25)


def foot(s, text):
    tf = box(s, L, FOOT_Y, Inches(9.0), Inches(0.3))
    para(tf, text, T_NOTE, color=LIGHT, first=True, after=0)


def card(s, l, t, w, h, accent, title, body):
    """The author's card: a wash panel with a narrow accent bar down its left."""
    rect(s, l, t, w, h, WASH)
    rect(s, l, t, Inches(0.07), h, accent)
    tf = box(s, l + Inches(0.25), t + Inches(0.14), w - Inches(0.45), h - Inches(0.28))
    para(tf, title, T_CARD, bold=True, first=True, after=4)
    para(tf, body, T_BODY, color=GREY, after=0, line=1.25)


def stat(s, t, figure, label, color=RED):
    tf = box(s, RIGHT_X, t, RIGHT_W, Inches(0.55))
    para(tf, figure, T_TITLE, bold=True, color=color, first=True, after=0, line=1.0)
    tf = box(s, RIGHT_X, t + Inches(0.52), RIGHT_W, Inches(0.75))
    para(tf, label, T_BODY, color=GREY, first=True, after=0, line=1.25)


def picture(s, name, l, t, w):
    pic = s.shapes.add_picture(SHOT + name, l, t, width=w)
    rect(s, l, t, pic.width, pic.height, None, RULE)
    return pic


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


# ============================================================ 1. Title
s = slide()
rect(s, 0, 0, W, H, WHITE)
s.shapes.add_picture(LOGO, Inches(0.38), Inches(1.55), height=Inches(0.95))
tf = box(s, Inches(0.38), Inches(2.80), Inches(11.5), Inches(1.6))
para(tf, "Preparing the People Who Teach", 31, bold=True, first=True, after=6, line=1.12)
para(tf, "A Casual Academic Management System, and the training to go with it",
     21, color=GREY, after=0, line=1.2)
rect(s, Inches(0.38), Inches(4.62), Inches(1.5), Pt(2.5), RED)
tf = box(s, Inches(0.38), Inches(5.00), Inches(11.5), Inches(1.0))
para(tf, "Wei Song  ·  School of Information and Communication Technology", 21,
     first=True, after=6)
para(tf, "Proposal to the Head of School  —  September 2026", 14, color=GREY, after=0)
notes(s, "Two halves. A system that makes tutors findable — built, deployed, and I can "
         "demonstrate it today. And the training program, which is what I am asking the "
         "School to endorse.")

# ============================================================ 2. The problem
s = slide()
y = head(s, "We Appoint Tutors Every Trimester,\nand Prepare None of Them",
         "A new tutor's readiness depends entirely on which convenor hired them.")
tf = box(s, L, y, Inches(7.1), Inches(3.4))
bullets(tf, [
    "Some receive a thorough handover; some are told the room number",
    "No School standard for what a tutor should be able to do before teaching",
    "No shared model of how an ICT tutorial or lab is run",
    "No marking calibration between tutors on the same course",
    "No record of who has been prepared, so nothing carries forward",
])
stat(s, Inches(2.00), "187", "Courses in the School, staffed by casual academics")
stat(s, Inches(3.65), "Every trimester", "We appoint tutors, demonstrators and markers", BLUE)
stat(s, Inches(5.30), "No standard", "For teaching preparation, at any level of the School", ORANGE)
foot(s, "The gap is not the convenors' — there is nothing for them to hand over to.")
notes(s, "This is not a criticism of convenors. There is no shared model, no materials "
         "and no record of who has taught what, so each of them starts from nothing.")

# ============================================================ 3. The gap
s = slide()
y = head(s, "Griffith Covers Compliance.\nNobody Covers Teaching.",
         "What a casual academic receives today, and what it does not include.")
rows = [("Mandatory onboarding", "WHS, integrity, privacy, equity"),
        ("Casual Staff Time Recording", "Timesheets and payment"),
        ("How to engage sessional staff", "Guidance for the hiring academic"),
        ("Tutoring for Success", "A different program, a different cohort")]
colw = [Inches(4.6), Inches(5.4), Inches(2.6)]
rect(s, L, y, sum(colw, Emu(0)), Inches(0.46), INK)
cx = L
for i, h2 in enumerate(["Provision", "What it covers", "Teaching?"]):
    tf = box(s, cx + Inches(0.22), y + Inches(0.12), colw[i] - Inches(0.34), Inches(0.3))
    para(tf, h2, T_BODY, bold=True, color=WHITE, first=True, after=0)
    cx += colw[i]
ry = y + Inches(0.46)
for r, (a, b) in enumerate(rows):
    rect(s, L, ry, sum(colw, Emu(0)), Inches(0.62), WHITE if r % 2 else WASH)
    cx = L
    for i, txt in enumerate((a, b, "No")):
        tf = box(s, cx + Inches(0.22), ry + Inches(0.16), colw[i] - Inches(0.34), Inches(0.4))
        para(tf, txt, T_BODY, bold=(i != 1),
             color=INK if i == 0 else (RED if i == 2 else GREY), first=True, after=0)
        cx += colw[i]
    ry += Inches(0.62)
tf = box(s, L, ry + Inches(0.45), CONTENT, Inches(1.1))
para(tf, "Compliance and payroll, done well. Nothing that prepares somebody to teach.",
     T_LEAD, bold=True, first=True, after=8, line=1.25)
para(tf, "The gap sits at School level: how a computing lab actually runs is not "
         "something a central program can cover.", T_BODY, color=GREY, after=0, line=1.3)
foot(s, "Griffith public sources, September 2026  ·  School Manager to confirm nothing "
        "internal already exists before this proceeds")
notes(s, "Checked against public Griffith sources only — I cannot see the staff intranet. "
         "The written proposal says we should confirm with the School Manager first. If "
         "something internal exists, we adopt it rather than duplicate it.")

# ============================================================ 4. The cycle
s = slide()
y = head(s, "One Cycle, Every Trimester,\nwith a Named Owner",
         "Steps 1, 2, 3 and 5 already happen — just not consistently, and nowhere on record.")
steps = [("1", "Invite", "Convenors invite students\nwho did well", "Convenor", BLUE),
         ("2", "Apply", "One standing application:\nexperience, ranked courses", "Candidate", BLUE),
         ("3", "Select", "Search by course.\nContact directly", "Convenor", BLUE),
         ("4", "Prepare", "Week 0 workshop\nThree hours", "Coordinator", RED),
         ("5", "Teach", "Shared tutorial model.\nWeek 3 check-in", "Convenor", BLUE)]
cw, gap = Inches(2.40), Inches(0.15)
x0 = L
for i, (n, title, bodytxt, owner, accent) in enumerate(steps):
    x = x0 + i * (cw + gap)
    hi = accent is RED
    rect(s, x, y, cw, Inches(2.85), WASH)
    rect(s, x, y, cw, Inches(0.07), accent)
    tf = box(s, x + Inches(0.24), y + Inches(0.28), cw - Inches(0.48), Inches(0.34))
    para(tf, f"STEP {n}", T_NOTE, bold=True, color=accent, first=True, after=6)
    tf = box(s, x + Inches(0.24), y + Inches(0.68), cw - Inches(0.48), Inches(0.4))
    para(tf, title, T_CARD + 2, bold=True, color=RED if hi else INK, first=True, after=8)
    tf = box(s, x + Inches(0.24), y + Inches(1.20), cw - Inches(0.48), Inches(1.0))
    lines = bodytxt.split("\n")
    for j, ln in enumerate(lines):
        para(tf, ln, T_BODY, color=GREY, first=(j == 0), after=2, line=1.25)
    tf = box(s, x + Inches(0.24), y + Inches(2.35), cw - Inches(0.48), Inches(0.3))
    para(tf, owner.upper(), T_NOTE, bold=True, color=LIGHT, first=True, after=0)
tf = box(s, L, y + Inches(3.20), CONTENT, Inches(0.9))
para(tf, "▪  Step 4 is the new part — and the invitation in Step 1 becomes proactive "
         "and merit-based, rather than students finding out by knowing somebody.",
     T_LEAD, first=True, after=0, line=1.3)
notes(s, "Two things change. Strong students get asked. And training becomes a School "
         "standard rather than a convenor's discretion.")

# ============================================================ 5. The system
s = slide()
y = head(s, "Half of It Already Exists",
         "The Casual Academic Management System — built, deployed, and in use today.")
picture(s, "00-entrance-candidate.png", Inches(0.9), y, Inches(5.4))
picture(s, "01-entrance-staff.png", Inches(7.0), y, Inches(5.4))
tf = box(s, Inches(0.9), y + Inches(3.55), Inches(5.4), Inches(0.4))
para(tf, "Candidates have their own link", T_BODY, bold=True, color=GREY,
     first=True, align=PP_ALIGN.CENTER, after=0)
tf = box(s, Inches(7.0), y + Inches(3.55), Inches(5.4), Inches(0.4))
para(tf, "Convenors have theirs", T_BODY, bold=True, color=GREY,
     first=True, align=PP_ALIGN.CENTER, after=0)
tf = box(s, L, y + Inches(4.15), CONTENT, Inches(0.7))
para(tf, "▪  One place where senior students put themselves forward, and convenors "
         "find out who can teach their course.", T_LEAD, first=True, after=0, line=1.3)
notes(s, "Two links. One goes to students, one to staff. Each audience sees only what "
         "applies to them. Built at no cost to the School.")

# ================================================= 6-9. The three views
views = [
    ("The Candidate's View", "One page, one form — everything a student needs to do.",
     "10-candidate.png",
     ["Details, teaching experience, ranked course choices, a supporting statement",
      "Always open — no recruitment round to wait for",
      "Starts blank each visit, so nothing stale is resubmitted by accident"],
     "Students have exactly one job here, so they get exactly one page and no navigation "
     "to learn."),
    ("The Convenor's View", "Every candidate in the School, searchable by course.",
     "11-convenor.png",
     ["Everyone registered — not only this trimester's applicants",
      "Two columns: what they have taught, and what they have applied for",
      "Filter by a course; those who have taught it are listed first"],
     "This is the page that addresses the real problem. Convenors do not struggle to run "
     "a selection process — they struggle to find out who is available at all. Note the "
     "people who have taught a course but did not apply this time; they were invisible "
     "before."),
    ("One Candidate", "Everything known about them, in one place.",
     "21-convenor-detail.png",
     ["Full teaching history, by course and trimester",
      "Their applied courses, in their own ranked order",
      "Their statement, contact details and availability"],
     "No approval workflow and no shortlisting buttons. The system informs the convenor; "
     "the decision and the conversation stay theirs."),
    ("The Administrator's View", "The School-wide picture, and the accounts behind it.",
     "12-admin.png",
     ["Every candidate, across all 187 courses",
      "Create and manage convenor accounts",
      "Export to CSV at any point, for the School office"],
     "As Coordinator this is where I would identify first-time tutors for the Week 0 "
     "workshop — derived from the records, rather than by asking around."),
]
for title, lead, img, pts, note in views:
    s = slide()
    y = head(s, title, lead)
    picture(s, img, L, y, Inches(8.1))
    bx = L + Inches(8.45)
    bw = W - bx - L
    for i, b in enumerate(pts):
        ty = y + Inches(0.18) + i * Inches(1.25)
        rect(s, bx, ty, Inches(0.07), Inches(0.85), BLUE)
        tf = box(s, bx + Inches(0.25), ty - Inches(0.02), bw - Inches(0.25), Inches(1.1))
        para(tf, b, T_BODY, color=GREY, first=True, after=0, line=1.35)
    notes(s, note)

# ============================================================ 10. Workshop
s = slide()
y = head(s, "The Missing Half: A Week 0 Workshop",
         "Three hours, before teaching starts, for everyone tutoring for the first time.")
seg = [("What a tutor does", "Not a second lecturer. Find out what students did not understand."),
       ("The tutorial model", "Review · supported practice · consolidation. Draft your own first session."),
       ("Micro-teaching", "Deliver a five-minute concept review to a peer, with feedback both ways."),
       ("Helping without answering", "Hands off the keyboard. Ask before you tell."),
       ("Marking consistently", "Everyone marks the same submission, then compares the spread."),
       ("Conduct and escalation", "Boundaries, integrity, students in difficulty. Escalate, don't absorb."),
       ("Senior tutor panel", "Experienced ICT tutors — the most credible voice in the room.")]
ty = y
for i, (h2, sub) in enumerate(seg):
    accent = RED if i in (1, 3) else BLUE
    rect(s, L, ty + Inches(0.04), Inches(0.07), Inches(0.42), accent)
    tf = box(s, L + Inches(0.28), ty, Inches(11.9), Inches(0.5))
    para(tf, h2, T_CARD, bold=True, first=True, after=2)
    para(tf, sub, T_BODY, color=GREY, after=0, line=1.2)
    ty += Inches(0.62)
tf = box(s, L, ty + Inches(0.2), CONTENT, Inches(0.6))
para(tf, "▪  Half the time is spent doing, not listening.", T_LEAD, bold=True,
     color=RED, first=True, after=0)
foot(s, "Deliberately excludes University compliance training and course-specific "
        "content — both are covered elsewhere.")
notes(s, "Not compliance training: Griffith does that already and repeating it would "
         "waste the only three hours the School gets. Not course-specific either — "
         "convenors brief their own tutors. This is what is common to all of them.")

# ============================================================ 11. Evaluation
s = slide()
y = head(s, "How We Would Know It Worked",
         "Measures chosen so that a failed pilot shows up as failed.")
cards = [("Coverage", "90% of first-time tutors trained before they teach", BLUE),
         ("Convenor judgement", "Better prepared than previous trimesters? This decides continuation.", RED),
         ("Participant judgement", "“Prepared to run my first tutorial” — 4.0 out of 5", BLUE),
         ("Retention", "60% of tutors return the following trimester", GREEN)]
cw2, gap2 = Inches(6.15), Inches(0.30)
for i, (t, b, c) in enumerate(cards):
    x = L + (i % 2) * (cw2 + gap2)
    ty = y + (i // 2) * Inches(1.45)
    card(s, x, ty, cw2, Inches(1.25), c, t, b)
rect(s, L, y + Inches(3.15), Inches(12.60), Inches(1.30), WASH)
rect(s, L, y + Inches(3.15), Inches(0.07), Inches(1.30), RED)
tf = box(s, L + Inches(0.30), y + Inches(3.35), Inches(12.0), Inches(1.0))
para(tf, "What would count as failure", T_CARD, bold=True, color=RED, first=True, after=7)
para(tf, "Convenors bypass the system  ·  No difference in preparedness  ·  "
         "First-time tutors do not attend  ·  It cannot run without me",
     T_BODY, color=GREY, after=0, line=1.3)
notes(s, "A program that cannot fail its own evaluation is not being evaluated. The last "
         "condition matters most: if it depends on one person, it is a hobby with a "
         "budget rather than a School process.")

# ============================================================ 12. The ask
s = slide()
y = head(s, "What I Am Asking For", "Four decisions.")
asks = [("Endorse the cycle", "Adopt the trimester process as School practice"),
        ("Endorse the expectation", "First-time tutors attend Week 0 before they teach"),
        ("Confirm the owner", "Casual Academic Coordinator — proposed: Wei Song"),
        ("Agree a pilot", "Run it once, review against the measures, then decide")]
ty = y
for i, (t, b) in enumerate(asks):
    rect(s, L, ty, Inches(12.60), Inches(0.95), WASH)
    rect(s, L, ty, Inches(0.07), Inches(0.95), RED)
    tf = box(s, L + Inches(0.35), ty + Inches(0.24), Inches(0.6), Inches(0.5))
    para(tf, str(i + 1), 22, bold=True, color=RED, first=True, after=0)
    tf = box(s, L + Inches(1.05), ty + Inches(0.18), Inches(11.0), Inches(0.65))
    para(tf, t, T_CARD, bold=True, first=True, after=3)
    para(tf, b, T_BODY, color=GREY, after=0)
    ty += Inches(1.08)
tf = box(s, L, ty + Inches(0.30), CONTENT, Inches(0.8))
para(tf, "▪  The system is built and the materials are written. What is needed is "
         "one trimester to try it.", T_LEAD, bold=True, first=True, after=0, line=1.3)
notes(s, "The pilot framing is deliberate. I am not asking for a permanent commitment to "
         "something untested — one trimester, measured against criteria agreed in "
         "advance.")

# ============================================================ 13. Try it
s = slide()
y = head(s, "Please Try It Yourself",
         "The system is live. These accounts are for you to use during or after this meeting.")

cw3, gap3 = Inches(6.15), Inches(0.30)

rect(s, L, y, cw3, Inches(2.35), WASH)
rect(s, L, y, Inches(0.07), Inches(2.35), BLUE)
tf = box(s, L + Inches(0.30), y + Inches(0.22), cw3 - Inches(0.55), Inches(1.95))
para(tf, "Candidate", T_CARD, bold=True, color=BLUE, first=True, after=10)
para(tf, "wweisong.github.io/griffith-ict-casual-academic", T_BODY, bold=True, after=12)
para(tf, "sample.liam.chen@griffithuni.edu.au", T_CARD, bold=True, after=4)
para(tf, "SamplePortal#2027", T_CARD, bold=True, color=RED, after=8)
para(tf, "Opens the application form a student fills in.", T_BODY, color=GREY, after=0)

rx = L + cw3 + gap3
rect(s, rx, y, cw3, Inches(2.35), WASH)
rect(s, rx, y, Inches(0.07), Inches(2.35), RED)
tf = box(s, rx + Inches(0.30), y + Inches(0.22), cw3 - Inches(0.55), Inches(1.95))
para(tf, "Course convenor", T_CARD, bold=True, color=RED, first=True, after=10)
para(tf, "wweisong.github.io/griffith-ict-casual-academic/#/staff", T_BODY, bold=True, after=12)
para(tf, "sample.convenor@griffith.edu.au", T_CARD, bold=True, after=4)
para(tf, "SamplePortal#2027", T_CARD, bold=True, color=RED, after=8)
para(tf, "Opens the candidate list — the page that matters most.", T_BODY, color=GREY, after=0)

picture(s, "22-admin-detail.png", Inches(4.55), y + Inches(2.55), Inches(4.25))
tf = box(s, L, y + Inches(5.25), CONTENT, Inches(0.32))
para(tf, "Administrator access on request. The candidates shown are sample records, "
         "not real students.", T_NOTE, color=LIGHT, first=True, align=PP_ALIGN.CENTER, after=0)
notes(s, "Please do open it. The convenor account shows the page that matters most. "
         "Everything — proposal, process, workshop run sheet, code of conduct — is at "
         "github.com/WWEISONG/griffith-ict-casual-academic")

# ============================================================ 14. Close
s = slide()
rect(s, 0, 0, W, H, WHITE)
s.shapes.add_picture(LOGO, Inches(0.38), Inches(1.75), height=Inches(0.85))
tf = box(s, Inches(0.38), Inches(2.95), Inches(11.6), Inches(1.8))
para(tf, "Students Remember Their Tutors", 31, bold=True, first=True, after=10, line=1.12)
para(tf, "They are the staff our students see most, and the only ones we have never "
         "prepared.", 21, color=GREY, after=0, line=1.25)
rect(s, Inches(0.38), Inches(4.85), Inches(1.5), Pt(2.5), RED)
tf = box(s, Inches(0.38), Inches(5.22), Inches(11.5), Inches(1.0))
para(tf, "Wei Song  ·  w.song@griffith.edu.au", 17, bold=True, first=True, after=5)
para(tf, "github.com/WWEISONG/griffith-ict-casual-academic", 14, color=GREY, after=0)
notes(s, "Everything is written up in the repository — the gap analysis with sources, "
         "the full process, the workshop run sheet, the code of conduct and the "
         "evaluation plan.")

out = "docs/tutor-program/HoS-briefing.pptx"
prs.save(out)
print(f"{out}  ·  {len(prs.slides._sldIdLst)} slides")
